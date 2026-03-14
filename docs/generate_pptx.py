"""
IP Assist — 고객사 전달용 통합 문서 (PowerPoint)
시스템구성도 + 상세스펙 + 이용자 매뉴얼 + QA 체크리스트
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── 색상 상수 ──
TEAL = RGBColor(0x0D, 0x94, 0x88)
DARK = RGBColor(0x22, 0x22, 0x22)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_BG = RGBColor(0xF5, 0xF5, 0xF5)
ORANGE = RGBColor(0xEA, 0x58, 0x0C)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)
RED = RGBColor(0xDC, 0x26, 0x26)
GREEN = RGBColor(0x16, 0xA3, 0x4A)
HEADER_BG = RGBColor(0x33, 0x33, 0x33)
TH_BG = RGBColor(0xE8, 0xE8, 0xE8)
TEAL_LIGHT = RGBColor(0xF0, 0xFD, 0xFA)
ORANGE_LIGHT = RGBColor(0xFF, 0xF7, 0xED)
PURPLE_LIGHT = RGBColor(0xF5, 0xF3, 0xFF)


def add_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, w, h, fill_color=None, border_color=None, border_width=Pt(1), text="", font_size=10, font_color=DARK, bold=False, align=PP_ALIGN.CENTER, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE):
    shape = slide.shapes.add_shape(shape_type, left, top, w, h)
    shape.shadow.inherit = False
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.color.rgb = font_color
        run.font.bold = bold
        run.font.name = "맑은 고딕"
    return shape


def set_cell(cell, text, font_size=9, bold=False, color=DARK, fill=None, align=PP_ALIGN.LEFT):
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "맑은 고딕"
    if fill:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def add_table(slide, rows_data, left, top, width, col_widths, row_height=Inches(0.38), header_bg=HEADER_BG, header_color=WHITE):
    """rows_data: list of lists. First row = header."""
    n_rows = len(rows_data)
    n_cols = len(rows_data[0])
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, Inches(row_height / Inches(1) * n_rows))
    table = table_shape.table
    for ci, cw in enumerate(col_widths):
        table.columns[ci].width = cw
    for ri, row in enumerate(rows_data):
        for ci, val in enumerate(row):
            cell = table.cell(ri, ci)
            if ri == 0:
                set_cell(cell, val, font_size=10, bold=True, color=header_color, fill=header_bg, align=PP_ALIGN.CENTER)
            else:
                is_first = ci == 0
                set_cell(cell, val, font_size=9, bold=is_first, color=DARK, fill=TH_BG if is_first else WHITE, align=PP_ALIGN.CENTER if is_first else PP_ALIGN.LEFT)
    return table_shape


def add_title_bar(slide, title, subtitle=""):
    # 상단 바
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.0))
    bar.fill.solid()
    bar.fill.fore_color.rgb = HEADER_BG
    bar.line.fill.background()
    bar.shadow.inherit = False
    tf = bar.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.space_before = Pt(6)
    run = p.add_run()
    run.text = f"    {title}"
    run.font.size = Pt(24)
    run.font.color.rgb = WHITE
    run.font.bold = True
    run.font.name = "맑은 고딕"
    if subtitle:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run()
        r2.text = f"    {subtitle}"
        r2.font.size = Pt(12)
        r2.font.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)
        r2.font.name = "맑은 고딕"
    # 하단 틸 라인
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.0), prs.slide_width, Pt(4))
    line.fill.solid()
    line.fill.fore_color.rgb = TEAL
    line.line.fill.background()
    line.shadow.inherit = False


def add_footer(slide, text="업무마켓9 — IP Assist | (주)캐스팅엔 | 2026.03"):
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(7.1), Inches(12), Inches(0.3))
    p = tx.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = text
    r.font.size = Pt(8)
    r.font.color.rgb = GRAY
    r.font.name = "맑은 고딕"


def add_text_box(slide, left, top, w, h, text, font_size=11, color=DARK, bold=False, align=PP_ALIGN.LEFT):
    tx = slide.shapes.add_textbox(left, top, w, h)
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(font_size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.name = "맑은 고딕"
    return tx


# ═══════════════════════════════════════════════════
# SLIDE 1 — 표지
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, HEADER_BG)

add_text_box(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(0.5),
             "간접구매 AI 상담도우미", font_size=16, color=RGBColor(0xAA, 0xAA, 0xAA), align=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(2.2), Inches(10), Inches(1.2),
             "시스템 구현 및 설계 문서", font_size=36, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# 틸 바
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(3.5), Inches(4.3), Pt(4))
bar.fill.solid()
bar.fill.fore_color.rgb = TEAL
bar.line.fill.background()
bar.shadow.inherit = False

add_text_box(slide, Inches(1.5), Inches(4.0), Inches(10), Inches(0.5),
             "IP Assist (업무마켓9)", font_size=20, color=TEAL, bold=True, align=PP_ALIGN.CENTER)

info_lines = "프로젝트명: 간접구매 AI 상담도우미 구축\n문서 버전: v1.0  |  작성일: 2026. 03. 08\n작성: (주)캐스팅엔 개발팀"
add_text_box(slide, Inches(3.5), Inches(5.2), Inches(6), Inches(1.2),
             info_lines, font_size=13, color=RGBColor(0xBB, 0xBB, 0xBB), align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════
# SLIDE 2 — 목차
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "목차", "Contents")
add_footer(slide)

toc_items = [
    ("01", "시스템 구성도", "외부망/내부망 아키텍처, 기술 스택"),
    ("02", "기술 스택 상세", "Gemini AI, Embedding 1536, pgvector, FastAPI"),
    ("03", "멀티에이전트 오케스트레이션", "7개 전문 에이전트 역할 및 병렬 처리"),
    ("04", "RAG 파이프라인", "벡터검색 + BM25 + RRF 리랭킹"),
    ("05", "데이터베이스 설계", "8개 테이블 스키마, HNSW 벡터 인덱스"),
    ("06", "이용자 가이드 — 챗봇 상담", "자연어 질의 → AI 실시간 스트리밍 답변"),
    ("07", "이용자 가이드 — RFP 작성", "4단계 페이즈 전환, 9개 RFP 유형"),
    ("08", "이용자 가이드 — 공급업체 추천", "카테고리 매칭, PDF 다운로드"),
    ("09", "관리자 포털 가이드", "지식베이스, 헌법, 사용자 관리"),
    ("10", "QA 테스트 체크리스트", "68개 점검 항목, 판정 기준"),
]

y_start = Inches(1.4)
for i, (num, title, desc) in enumerate(toc_items):
    y = y_start + Inches(i * 0.55)
    # 번호
    add_shape(slide, Inches(1.0), y, Inches(0.6), Inches(0.42), fill_color=TEAL, text=num, font_size=14, font_color=WHITE, bold=True)
    # 제목
    add_text_box(slide, Inches(1.8), y, Inches(4), Inches(0.42), title, font_size=14, color=DARK, bold=True)
    # 설명
    add_text_box(slide, Inches(6.0), y + Pt(2), Inches(6), Inches(0.38), desc, font_size=10, color=GRAY)
    # 구분선
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), y + Inches(0.46), Inches(11.3), Pt(0.5))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    line.line.fill.background()
    line.shadow.inherit = False

# ═══════════════════════════════════════════════════
# SLIDE 3 — 시스템 구성도
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "01  시스템 구성도", "System Architecture")
add_footer(slide)

# 외부망 영역
add_shape(slide, Inches(0.4), Inches(1.4), Inches(2.0), Inches(5.2),
          fill_color=RGBColor(0xFF, 0xF8, 0xE1), border_color=HEADER_BG, border_width=Pt(1.5),
          shape_type=MSO_SHAPE.RECTANGLE)
add_text_box(slide, Inches(0.5), Inches(1.5), Inches(1.8), Inches(0.3),
             "외부망", font_size=12, color=DARK, bold=True, align=PP_ALIGN.CENTER)

# 외부망 박스들
add_shape(slide, Inches(0.6), Inches(2.1), Inches(1.6), Inches(0.7),
          fill_color=WHITE, border_color=GRAY, text="사용자\n(웹 브라우저)", font_size=9)
add_shape(slide, Inches(0.6), Inches(3.1), Inches(1.6), Inches(0.7),
          fill_color=WHITE, border_color=GRAY, text="관리자\n(관리 포털)", font_size=9)
add_shape(slide, Inches(0.6), Inches(4.1), Inches(1.6), Inches(0.7),
          fill_color=WHITE, border_color=GRAY, text="Google Gemini\nAPI", font_size=9)

# 내부망 영역
add_shape(slide, Inches(2.7), Inches(1.4), Inches(10.2), Inches(5.2),
          fill_color=LIGHT_BG, border_color=HEADER_BG, border_width=Pt(1.5),
          shape_type=MSO_SHAPE.RECTANGLE)
add_text_box(slide, Inches(2.8), Inches(1.5), Inches(2), Inches(0.3),
             "내부망 (클라우드)", font_size=12, color=DARK, bold=True)

# Firebase Hosting
add_shape(slide, Inches(3.1), Inches(2.1), Inches(1.7), Inches(0.9),
          fill_color=TEAL_LIGHT, border_color=TEAL, border_width=Pt(2),
          text="Firebase Hosting\nCDN / SSL", font_size=9, font_color=TEAL, bold=True)

# 화살표 텍스트
add_text_box(slide, Inches(4.9), Inches(2.3), Inches(0.5), Inches(0.3), "→", font_size=18, color=GRAY, align=PP_ALIGN.CENTER)

# Frontend SPA
add_shape(slide, Inches(5.3), Inches(2.1), Inches(1.7), Inches(0.9),
          fill_color=TEAL_LIGHT, border_color=TEAL, border_width=Pt(2),
          text="Frontend (SPA)\nReact 19 + Vite 7", font_size=9, font_color=TEAL, bold=True)

# 화살표
add_text_box(slide, Inches(7.1), Inches(2.3), Inches(0.8), Inches(0.5), "HTTPS\nSSE →", font_size=8, color=GRAY, align=PP_ALIGN.CENTER)

# FastAPI Backend (큰 박스)
add_shape(slide, Inches(7.9), Inches(1.9), Inches(3.0), Inches(2.8),
          fill_color=WHITE, border_color=TEAL, border_width=Pt(2),
          shape_type=MSO_SHAPE.RECTANGLE)
add_text_box(slide, Inches(8.0), Inches(2.0), Inches(2.8), Inches(0.3),
             "FastAPI Backend (Cloud Run)", font_size=10, color=TEAL, bold=True, align=PP_ALIGN.CENTER)

# API 모듈들
modules = ["Chat", "RFP", "Knowledge", "Suppliers", "Admin"]
for i, m in enumerate(modules):
    x = Inches(8.05) + Inches(i * 0.56)
    add_shape(slide, x, Inches(2.45), Inches(0.52), Inches(0.32),
              fill_color=TEAL_LIGHT, border_color=TEAL,
              text=m, font_size=7, font_color=TEAL, bold=True)

# 에이전트 모듈
agents = ["Orchestrator", "RAG Engine", "Constitution"]
for i, a in enumerate(agents):
    x = Inches(8.15) + Inches(i * 0.92)
    add_shape(slide, x, Inches(3.0), Inches(0.85), Inches(0.32),
              fill_color=ORANGE_LIGHT, border_color=ORANGE,
              text=a, font_size=7, font_color=ORANGE, bold=True)

# 검색 엔진 표시
add_text_box(slide, Inches(8.0), Inches(3.5), Inches(2.8), Inches(0.25),
             "Vector + BM25 + RRF 리랭킹", font_size=8, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(8.0), Inches(3.85), Inches(2.8), Inches(0.25),
             "멀티에이전트 (7개 병렬 오케스트레이션)", font_size=8, color=ORANGE, align=PP_ALIGN.CENTER)

# 화살표 → DB
add_text_box(slide, Inches(11.0), Inches(2.8), Inches(0.5), Inches(0.3), "→", font_size=18, color=GRAY, align=PP_ALIGN.CENTER)

# Supabase DB
add_shape(slide, Inches(11.4), Inches(2.2), Inches(1.3), Inches(1.2),
          fill_color=PURPLE_LIGHT, border_color=PURPLE, border_width=Pt(2),
          text="Supabase\nPostgreSQL\npgvector\n1536차원", font_size=9, font_color=PURPLE, bold=True)

# AI 서비스 영역
add_shape(slide, Inches(3.1), Inches(3.5), Inches(4.5), Inches(1.5),
          fill_color=ORANGE_LIGHT, border_color=ORANGE, border_width=Pt(1.5),
          shape_type=MSO_SHAPE.RECTANGLE)
add_text_box(slide, Inches(3.2), Inches(3.55), Inches(4.3), Inches(0.25),
             "AI / ML 서비스 (Google Gemini)", font_size=10, color=ORANGE, bold=True)

ai_boxes = [
    ("Gemini 2.5 Flash\nRAG 답변 생성\n스트리밍", Inches(3.3)),
    ("Gemini Flash-Lite\n데이터 정제\n의도 분류", Inches(4.8)),
    ("Gemini Embedding\n문서 벡터화\n1536차원", Inches(6.3)),
]
for text, x in ai_boxes:
    add_shape(slide, x, Inches(3.9), Inches(1.3), Inches(0.9),
              fill_color=WHITE, border_color=ORANGE, text=text, font_size=8)

# 인프라 행
infra = [
    ("Cloud Run\nasia-northeast3\nmin-instances: 1", Inches(3.1)),
    ("Firebase Hosting\nCDN 글로벌 배포", Inches(5.0)),
    ("Supabase\nap-northeast-2\n서울 리전", Inches(6.9)),
    ("Docker\nPython 3.11-slim\n컨테이너", Inches(8.8)),
]
for text, x in infra:
    add_shape(slide, x, Inches(5.3), Inches(1.6), Inches(0.9),
              fill_color=HEADER_BG, border_color=HEADER_BG, text=text, font_size=8, font_color=WHITE, bold=True)

# 범례
add_text_box(slide, Inches(10.8), Inches(5.4), Inches(2), Inches(0.8),
             "■ 프론트엔드/API\n■ AI/ML 서비스\n■ 데이터베이스\n■ 인프라", font_size=8, color=GRAY)

# ═══════════════════════════════════════════════════
# SLIDE 4 — 기술 스택 상세
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "02  기술 스택 상세", "Technology Stack Details")
add_footer(slide)

rows = [
    ["계층", "기술", "버전", "상세 설명"],
    ["Frontend", "React + Vite", "19.2 / 7.3", "SPA, React Router DOM 7, jsPDF PDF 생성, SSE 실시간 수신"],
    ["Backend", "Python FastAPI", "3.11 / 0.115", "Uvicorn ASGI, SSE 스트리밍, CORS, 멀티에이전트 오케스트레이션"],
    ["AI 생성", "Gemini 2.5 Flash", "최신", "RAG 답변 생성, 토큰 단위 스트리밍, 분류/제안/필드추출"],
    ["AI 정제", "Gemini 2.5 Flash-Lite", "최신", "PDF 데이터 정제, 의도 분류, 후속 질문 생성 (비용 효율)"],
    ["Embedding", "Gemini Embedding 001", "최신", "1536차원 벡터 임베딩, RETRIEVAL_DOCUMENT / RETRIEVAL_QUERY"],
    ["Vector DB", "PostgreSQL + pgvector", "15+", "Supabase 호스팅, HNSW 인덱스 (m=16, ef=64), 코사인 유사도"],
    ["검색", "rank-bm25 + RRF", "0.2.2", "BM25 키워드 검색, Reciprocal Rank Fusion 리랭킹 (k=60)"],
    ["Frontend 호스팅", "Firebase Hosting", "-", "글로벌 CDN, SSL, SPA 리라이트, JS/CSS 1년 캐시"],
    ["Backend 호스팅", "Google Cloud Run", "-", "asia-northeast3 (서울), min-instances=1, Docker 컨테이너"],
]
add_table(slide, rows, Inches(0.5), Inches(1.4), Inches(12.3),
          [Inches(1.5), Inches(2.3), Inches(1.2), Inches(7.3)], row_height=Inches(0.42))

# 핵심 스펙 하이라이트
specs = [
    ("임베딩 차원: 1536", TEAL),
    ("벡터 인덱스: HNSW", PURPLE),
    ("RRF k=60", ORANGE),
    ("FAQ Boost 1.5x", ORANGE),
    ("min-instances: 1", HEADER_BG),
    ("P95 응답: ~2.5s", TEAL),
]
for i, (text, color) in enumerate(specs):
    x = Inches(0.5) + Inches(i * 2.1)
    add_shape(slide, x, Inches(6.1), Inches(1.9), Inches(0.4),
              fill_color=color, text=text, font_size=10, font_color=WHITE, bold=True)

# ═══════════════════════════════════════════════════
# SLIDE 5 — 멀티에이전트 오케스트레이션
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "03  멀티에이전트 오케스트레이션", "Multi-Agent Orchestration System")
add_footer(slide)

rows = [
    ["#", "에이전트", "AI 모델", "역할", "실행"],
    ["1", "Orchestrator", "-", "에이전트 실행 순서 결정, 병렬 관리, SSE 스트림 조율", "항상"],
    ["2", "Classification", "Flash-Lite", "의도 분류 (구매/질문/RFP동의/필드입력), RFP 유형 추천", "병렬"],
    ["3", "Retrieval", "Embedding-001", "하이브리드 검색 (Vector+BM25+FAQ), RRF 리랭킹, 상위 5개", "병렬"],
    ["4", "Constitution", "Embedding-001", "사전검사(위반탐지), 규칙주입(컨텍스트보강), 사후검증", "순차"],
    ["5", "Generation", "Flash", "RAG 컨텍스트 기반 답변 생성, SSE 토큰 스트리밍", "순차"],
    ["6", "Suggestion", "Flash-Lite", "후속 질문 3~5개 생성 (비블로킹)", "비동기"],
    ["7", "RFP Extractor", "Flash-Lite", "대화에서 RFP 필드 자동 추출 (사업명, 예산, 기간 등)", "조건부"],
    ["8", "Prefetcher", "Embedding-001", "다음 질문 예측, 검색 결과 사전 캐싱 (5분 TTL)", "비동기"],
]
add_table(slide, rows, Inches(0.5), Inches(1.4), Inches(12.3),
          [Inches(0.4), Inches(1.5), Inches(1.4), Inches(6.5), Inches(0.8)], row_height=Inches(0.44))

# 실행 흐름 시각화
add_text_box(slide, Inches(0.5), Inches(5.6), Inches(12), Inches(0.3),
             "오케스트레이션 실행 순서", font_size=12, color=DARK, bold=True)

flow_steps = [
    ("① 헌법\n사전검사", HEADER_BG),
    ("② 페이즈\n트리거 감지", HEADER_BG),
    ("③ 병렬 실행\nClassification\n+ Retrieval", TEAL),
    ("④ 헌법\n규칙 주입", ORANGE),
    ("⑤ Generation\nSSE 스트리밍", TEAL),
    ("⑥ 사후검증\n(비동기)", GRAY),
    ("⑦ 후속질문\n생성", GRAY),
]
for i, (text, color) in enumerate(flow_steps):
    x = Inches(0.5) + Inches(i * 1.75)
    add_shape(slide, x, Inches(6.0), Inches(1.5), Inches(0.85),
              fill_color=color, text=text, font_size=8, font_color=WHITE, bold=True)
    if i < len(flow_steps) - 1:
        add_text_box(slide, x + Inches(1.5), Inches(6.25), Inches(0.3), Inches(0.3), "→", font_size=16, color=GRAY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════
# SLIDE 6 — RAG 파이프라인
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "04  RAG 파이프라인", "Retrieval-Augmented Generation Pipeline")
add_footer(slide)

# 왼쪽: 인제스션 파이프라인
add_text_box(slide, Inches(0.5), Inches(1.3), Inches(5), Inches(0.3),
             "문서 인제스션 파이프라인", font_size=13, color=DARK, bold=True)

ingest_steps = [
    ("PDF 업로드", "관리자가 PDF 파일 업로드"),
    ("텍스트 추출", "pdfplumber로 페이지별 텍스트 추출"),
    ("청크 분할", "400토큰 단위, 50토큰 오버랩"),
    ("AI 정제", "Gemini Flash-Lite로 노이즈 제거"),
    ("벡터 임베딩", "Gemini Embedding 001 → 1536차원"),
    ("DB 저장", "Supabase knowledge_chunks + HNSW 인덱스"),
]
for i, (title, desc) in enumerate(ingest_steps):
    y = Inches(1.8) + Inches(i * 0.72)
    add_shape(slide, Inches(0.5), y, Inches(0.45), Inches(0.45),
              fill_color=TEAL, text=str(i+1), font_size=14, font_color=WHITE, bold=True)
    add_text_box(slide, Inches(1.1), y, Inches(1.5), Inches(0.25), title, font_size=11, color=DARK, bold=True)
    add_text_box(slide, Inches(1.1), y + Inches(0.25), Inches(4.5), Inches(0.25), desc, font_size=9, color=GRAY)
    if i < len(ingest_steps) - 1:
        add_text_box(slide, Inches(0.6), y + Inches(0.5), Inches(0.3), Inches(0.2), "↓", font_size=12, color=TEAL, align=PP_ALIGN.CENTER)

# 오른쪽: 하이브리드 검색 + RRF
add_text_box(slide, Inches(6.5), Inches(1.3), Inches(6), Inches(0.3),
             "하이브리드 검색 + RRF 리랭킹", font_size=13, color=DARK, bold=True)

# 검색 소스 3개
search_sources = [
    ("FAQ 검색\nBoost: 1.5x\n임계값: 0.55", TEAL, Inches(6.5)),
    ("BM25 검색\nBoost: 1.2x\n키워드 매칭", ORANGE, Inches(8.5)),
    ("Vector 검색\nBoost: 1.0x\n임계값: 0.75", PURPLE, Inches(10.5)),
]
for text, color, x in search_sources:
    add_shape(slide, x, Inches(1.8), Inches(1.7), Inches(1.1),
              fill_color=color, text=text, font_size=9, font_color=WHITE, bold=True)

# 병렬 표시
add_text_box(slide, Inches(6.5), Inches(3.0), Inches(5.7), Inches(0.3),
             "▼  ThreadPoolExecutor 병렬 실행  ▼", font_size=10, color=TEAL, bold=True, align=PP_ALIGN.CENTER)

# RRF 공식
add_shape(slide, Inches(6.5), Inches(3.4), Inches(5.7), Inches(1.2),
          fill_color=RGBColor(0xF8, 0xF8, 0xF8), border_color=HEADER_BG, border_width=Pt(1.5),
          shape_type=MSO_SHAPE.RECTANGLE)
add_text_box(slide, Inches(6.7), Inches(3.5), Inches(5.3), Inches(0.3),
             "RRF (Reciprocal Rank Fusion)", font_size=12, color=DARK, bold=True, align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(6.7), Inches(3.9), Inches(5.3), Inches(0.3),
             "score(d) = Σ( boost_source / (k + rank(d)) )    k = 60", font_size=11, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(6.7), Inches(4.2), Inches(5.3), Inches(0.25),
             "3개 소스 결과 통합 → 중복 제거 → RRF 점수 계산 → 상위 5개 선정", font_size=9, color=GRAY, align=PP_ALIGN.CENTER)

# 성능 비교
add_text_box(slide, Inches(6.5), Inches(4.9), Inches(5.7), Inches(0.3),
             "검색 성능 (병렬 실행 효과)", font_size=11, color=DARK, bold=True)
perf_rows = [
    ["구분", "이전 (순차)", "현재 (병렬+RRF)", "개선"],
    ["FAQ 검색", "~100ms", "~100ms (병렬)", "-"],
    ["BM25 검색", "~150ms", "~150ms (병렬)", "-"],
    ["Vector 검색", "~200ms", "~200ms (병렬)", "-"],
    ["합계", "~450ms (순차 누적)", "~210ms (병렬+RRF)", "50%↓"],
]
add_table(slide, perf_rows, Inches(6.5), Inches(5.3), Inches(5.7),
          [Inches(1.2), Inches(1.5), Inches(1.5), Inches(1.0)], row_height=Inches(0.32))

# ═══════════════════════════════════════════════════
# SLIDE 7 — 데이터베이스 설계
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "05  데이터베이스 설계", "Supabase PostgreSQL + pgvector Schema")
add_footer(slide)

db_rows = [
    ["테이블명", "주요 컬럼", "인덱스", "용도"],
    ["knowledge_chunks", "category, doc_name, content, embedding(1536), metadata", "HNSW (embedding), B-tree (category, doc_name)", "RAG 벡터 검색 대상"],
    ["knowledge_faq", "category, question, answer, embedding(1536)", "HNSW (embedding)", "FAQ 우선 검색 (boost 1.5x)"],
    ["conversations", "session_id, messages(jsonb), rfp_data, rag_score, status", "-", "대화 이력 + RFP 데이터"],
    ["constitution_rules", "rule_type, content, embedding(1536), is_active", "-", "AI 가드레일 (헌법 규칙)"],
    ["suppliers", "name, category, score(0-20), match_rate(0-100), tags[]", "B-tree (category)", "공급업체 DB + 매칭"],
    ["admin_users", "name, email(UNIQUE), password_hash, role, department", "B-tree (email)", "관리자 인증 + 권한"],
    ["rfp_requests", "session_id, rfp_type, title, fields(jsonb), status", "-", "RFP 신청 관리"],
    ["rfp_templates", "type_key, name, fields(jsonb), sections[]", "-", "RFP 양식 정의"],
]
add_table(slide, db_rows, Inches(0.3), Inches(1.4), Inches(12.7),
          [Inches(1.8), Inches(4.5), Inches(3.2), Inches(2.5)], row_height=Inches(0.44))

# 벡터 검색 함수
add_text_box(slide, Inches(0.5), Inches(5.7), Inches(5), Inches(0.3),
             "벡터 검색 함수: match_chunks()", font_size=12, color=DARK, bold=True)
add_shape(slide, Inches(0.5), Inches(6.1), Inches(5.5), Inches(1.0),
          fill_color=RGBColor(0xF8, 0xF8, 0xF8), border_color=GRAY,
          text="match_chunks(query_embedding, match_count=5, category_filter=NULL)\n→ HNSW 인덱스 코사인 유사도 검색\n→ 1-(embedding <=> query) = similarity\n→ 카테고리 필터 + 상위 K개 반환",
          font_size=9, font_color=DARK, align=PP_ALIGN.LEFT, shape_type=MSO_SHAPE.RECTANGLE)

# 헌법 규칙 유형
add_text_box(slide, Inches(6.5), Inches(5.7), Inches(5), Inches(0.3),
             "AI 헌법 (가드레일) 규칙 6개", font_size=12, color=DARK, bold=True)
rules_rows = [
    ["유형", "규칙 내용", "검사"],
    ["거부조건", "특정 업체를 이유 없이 추천/선호 금지", "사전+사후"],
    ["행동제어", "근거 문서에 없는 조항 임의 추가 금지", "사후"],
    ["수치기준", "계약금액은 분류체계 문서 수치 인용 필수", "사후"],
    ["행동제어", "법적 판단 요구 시 전문가 상담 권고", "사전+사후"],
    ["거부조건", "개인정보 학습/기억 금지", "사전"],
    ["행동제어", "답변에 RAG 근거 문서 출처 명시 필수", "사후"],
]
add_table(slide, rules_rows, Inches(6.5), Inches(6.1), Inches(6.0),
          [Inches(1.0), Inches(3.8), Inches(1.0)], row_height=Inches(0.28))

# ═══════════════════════════════════════════════════
# SLIDE 8 — 이용자 가이드: 챗봇 상담
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "06  이용자 가이드 — 챗봇 상담", "User Guide: AI Chat Consultation")
add_footer(slide)

# 4단계 페이즈 플로우
phases = [
    ("Phase 1\nChat (상담)", "자유 질의응답\nRAG 기반 답변\n출처 표시", TEAL),
    ("Phase 2\nAsked (제안)", "구매 의도 감지\nRFP 작성 제안\n사용자 동의 대기", ORANGE),
    ("Phase 3\nFilling (작성)", "RFP 유형 선택\n필드 입력\n자동 추출", PURPLE),
    ("Phase 4\nComplete (완료)", "PDF 다운로드\n공급업체 추천\n세션 저장", HEADER_BG),
]
for i, (title, desc, color) in enumerate(phases):
    x = Inches(0.5) + Inches(i * 3.2)
    add_shape(slide, x, Inches(1.5), Inches(2.7), Inches(1.5),
              fill_color=color, text=f"{title}\n\n{desc}", font_size=10, font_color=WHITE, bold=True)
    if i < 3:
        add_text_box(slide, x + Inches(2.7), Inches(2.0), Inches(0.5), Inches(0.4), "→", font_size=24, color=TEAL, bold=True, align=PP_ALIGN.CENTER)

# 사용 방법
add_text_box(slide, Inches(0.5), Inches(3.3), Inches(12), Inches(0.3),
             "챗봇 사용 방법", font_size=13, color=DARK, bold=True)

steps = [
    ("1", "웹 브라우저 접속", "챗봇 URL 접속 → 왼쪽 채팅 영역 + 오른쪽 RFP 패널 표시"),
    ("2", "자연어 질문 입력", '하단 입력창에 질문 입력. 예: "교육 서비스 구매 절차가 어떻게 되나요?"'),
    ("3", "AI 실시간 답변 확인", "SSE 스트리밍으로 한 글자씩 답변 표시. 하단에 근거 문서 출처 명시."),
    ("4", "후속 질문 활용", "답변 후 3~5개 후속 질문 제안 버튼 표시. 클릭 시 해당 질문 자동 전송."),
    ("5", "구매 의도 자동 감지", 'AI가 구매 의도 인식 → "RFP를 작성해 드릴까요?" 자동 제안'),
]
for i, (num, title, desc) in enumerate(steps):
    y = Inches(3.8) + Inches(i * 0.6)
    add_shape(slide, Inches(0.5), y, Inches(0.4), Inches(0.4),
              fill_color=TEAL, text=num, font_size=12, font_color=WHITE, bold=True)
    add_text_box(slide, Inches(1.1), y, Inches(2.2), Inches(0.3), title, font_size=11, color=DARK, bold=True)
    add_text_box(slide, Inches(3.3), y + Pt(2), Inches(9.5), Inches(0.35), desc, font_size=10, color=GRAY)

# TIP 박스
add_shape(slide, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.45),
          fill_color=ORANGE_LIGHT, border_color=ORANGE, border_width=Pt(1),
          text='TIP: 구체적 질문일수록 정확한 답변!   "구매하고 싶어요" → "법정의무교육 위탁 운영 서비스를 200명 규모로 구매하고 싶습니다"',
          font_size=9, font_color=ORANGE, shape_type=MSO_SHAPE.RECTANGLE)

# ═══════════════════════════════════════════════════
# SLIDE 9 — 이용자 가이드: RFP 작성 + 공급업체 추천
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "07-08  RFP 작성 & 공급업체 추천", "RFP Generation Flow & Supplier Recommendation")
add_footer(slide)

# RFP 유형 9종
add_text_box(slide, Inches(0.5), Inches(1.3), Inches(5), Inches(0.3),
             "RFP 유형 (9종)", font_size=13, color=DARK, bold=True)

rfp_types = [
    "물품 구매 계약", "용역 계약", "수리/유지보수",
    "리스/임대", "공사/시설", "컨설팅/자문",
    "물품+유지보수", "임대+유지보수", "구매 또는 리스",
]
for i, t in enumerate(rfp_types):
    row, col = divmod(i, 3)
    x = Inches(0.5) + Inches(col * 2.0)
    y = Inches(1.8) + Inches(row * 0.5)
    color = TEAL if i == 0 else HEADER_BG
    add_shape(slide, x, y, Inches(1.8), Inches(0.38),
              fill_color=color, text=t, font_size=9, font_color=WHITE, bold=True)

# RFP 작성 방법
add_text_box(slide, Inches(0.5), Inches(3.5), Inches(5.5), Inches(0.3),
             "RFP 작성 방법 (2가지)", font_size=12, color=DARK, bold=True)

add_shape(slide, Inches(0.5), Inches(3.9), Inches(2.7), Inches(1.5),
          fill_color=TEAL_LIGHT, border_color=TEAL, border_width=Pt(2),
          text="방법 1 — 대화로 입력 (권장)\n\n채팅으로 자연어 입력\n→ AI가 필드 자동 추출\n→ RFP 패널에 자동 반영",
          font_size=9, font_color=DARK, shape_type=MSO_SHAPE.RECTANGLE)

add_shape(slide, Inches(3.5), Inches(3.9), Inches(2.7), Inches(1.5),
          fill_color=LIGHT_BG, border_color=GRAY, border_width=Pt(1),
          text="방법 2 — 직접 입력\n\n오른쪽 RFP 패널의\n각 필드를 클릭하여\n직접 값을 입력",
          font_size=9, font_color=DARK, shape_type=MSO_SHAPE.RECTANGLE)

# PDF 다운로드
add_text_box(slide, Inches(0.5), Inches(5.6), Inches(5.5), Inches(0.3),
             "PDF 다운로드", font_size=12, color=DARK, bold=True)
add_text_box(slide, Inches(0.5), Inches(5.95), Inches(5.5), Inches(0.7),
             "• 필수 필드 100% 입력 시 다운로드 가능\n• 나라장터 스타일 공식 문서 형태\n• 구성: 사업명, 발주기관, 담당자, 기간, 목표, 내용, 평가기준, 서명란",
             font_size=9, color=GRAY)

# 오른쪽: 공급업체 추천
add_text_box(slide, Inches(6.8), Inches(1.3), Inches(6), Inches(0.3),
             "공급업체 추천 프로세스", font_size=13, color=DARK, bold=True)

rec_steps = [
    ("1", "RFP 완성", "필수 필드 100% → Complete 페이즈 전환"),
    ("2", "카테고리 매칭", "RFP 카테고리 기반 공급업체 DB 검색"),
    ("3", "추천 목록 표시", "매칭률(%), 만족도 점수, 태그 정보"),
    ("4", "업체 선택", "추천 목록에서 업체 선택"),
    ("5", "RFP 전송", "선택한 업체에 제안요청서 전달"),
]
for i, (num, title, desc) in enumerate(rec_steps):
    y = Inches(1.8) + Inches(i * 0.65)
    add_shape(slide, Inches(6.8), y, Inches(0.4), Inches(0.4),
              fill_color=PURPLE, text=num, font_size=12, font_color=WHITE, bold=True)
    add_text_box(slide, Inches(7.4), y, Inches(1.8), Inches(0.25), title, font_size=11, color=DARK, bold=True)
    add_text_box(slide, Inches(7.4), y + Inches(0.25), Inches(5), Inches(0.25), desc, font_size=9, color=GRAY)

# 추천 알고리즘
add_text_box(slide, Inches(6.8), Inches(5.3), Inches(6), Inches(0.3),
             "공급업체 데이터 구조", font_size=12, color=DARK, bold=True)
sup_rows = [
    ["필드", "타입", "설명"],
    ["name", "text", "업체명"],
    ["category", "text", "간접구매 카테고리 (교육, IT, 시설 등)"],
    ["score", "0~20", "만족도 평점"],
    ["match_rate", "0~100%", "카테고리 매칭률"],
    ["tags", "array", "특성 태그 (ESG, 중소기업, 여성기업 등)"],
    ["status", "text", "상태 (active / review / inactive)"],
]
add_table(slide, sup_rows, Inches(6.8), Inches(5.65), Inches(5.7),
          [Inches(1.2), Inches(1.0), Inches(3.3)], row_height=Inches(0.26))

# ═══════════════════════════════════════════════════
# SLIDE 10 — 관리자 포털
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "09  관리자 포털 가이드", "Admin Portal Guide")
add_footer(slide)

admin_rows = [
    ["메뉴", "기능", "상세 설명"],
    ["대시보드", "시스템 현황", "지식베이스 청크 수, 대화 수, 공급업체 수, 헌법 조항 수 통계 카드 표시"],
    ["지식베이스", "PDF 업로드 + 현황", "PDF 업로드 → 자동 인제스션. 카테고리별 청크/FAQ/문서 수 확인. 10건 페이지네이션."],
    ["헌법 관리", "AI 가드레일 CRUD", "거부조건/행동제어/수치기준 3유형 규칙 추가/수정/삭제. 삭제 확인 모달. 자동 임베딩."],
    ["대화 이력", "세션별 대화 조회", "사용자명, 부서, 카테고리, RAG 점수, 상태 필터링. 세션 클릭 시 전체 대화 표시."],
    ["RFP 신청", "제출된 RFP 관리", "RFP 목록 조회, 상태 변경(제출/승인/반려), 삭제"],
    ["RFP 양식", "양식 편집기", "RFP 유형별 섹션/필드 구조 정의, 신규 생성/수정/삭제"],
    ["공급업체", "업체 DB 관리", "업체 등록/수정/삭제, CSV 일괄 등록, 카테고리/상태 필터링"],
    ["사용자 관리", "관리자 계정 CRUD", "이름/이메일/비밀번호/역할/부서. 역할: superadmin/admin/viewer. 활성화 토글."],
]
add_table(slide, admin_rows, Inches(0.3), Inches(1.4), Inches(12.7),
          [Inches(1.4), Inches(1.8), Inches(8.8)], row_height=Inches(0.48))

# 권한 체계
add_text_box(slide, Inches(0.5), Inches(6.0), Inches(3), Inches(0.3),
             "역할 기반 접근 제어", font_size=12, color=DARK, bold=True)
role_rows = [
    ["역할", "권한"],
    ["superadmin (최고관리자)", "전체 기능 접근, 사용자 관리 포함"],
    ["admin (관리자)", "데이터 관리 (사용자 관리 제한)"],
    ["viewer (뷰어)", "조회만 가능 (수정/삭제 불가)"],
]
add_table(slide, role_rows, Inches(0.5), Inches(6.35), Inches(5.5),
          [Inches(2.5), Inches(3.0)], row_height=Inches(0.30))

# 접속 정보
add_text_box(slide, Inches(7.0), Inches(6.0), Inches(3), Inches(0.3),
             "접속 정보", font_size=12, color=DARK, bold=True)
access_rows = [
    ["항목", "값"],
    ["챗봇 URL", "https://chatbot-copilot-82b0a.web.app"],
    ["관리자 URL", "위 URL + /admin"],
    ["관리자 계정", "admin@castingn.com"],
    ["인증 방식", "이메일 + SHA256 해시 비밀번호"],
]
add_table(slide, access_rows, Inches(7.0), Inches(6.35), Inches(5.5),
          [Inches(1.5), Inches(4.0)], row_height=Inches(0.26))

# ═══════════════════════════════════════════════════
# SLIDE 11 — QA 테스트 체크리스트 (1)
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "10  QA 테스트 체크리스트 (1/2)", "챗봇 기능 + 페이즈 전환 + RFP + AI 품질")
add_footer(slide)

qa1_rows = [
    ["#", "우선", "테스트 항목", "기대 결과", "P/F"],
    ["1", "HIGH", "챗봇 URL 접속", "페이지 정상 로딩, 채팅 입력창 표시", "☐"],
    ["2", "HIGH", '일반 질문 입력 ("간접구매란?")', "AI 실시간 스트리밍 답변 + 근거 출처 표시", "☐"],
    ["3", "HIGH", "SSE 스트리밍 동작", "텍스트가 한 글자씩 실시간 표시", "☐"],
    ["4", "MID", "후속 질문 제안 버튼", "답변 후 3~5개 제안 표시, 클릭 시 전송", "☐"],
    ["5", "MID", "연속 대화 5회 이상", "대화 이력 유지, 컨텍스트 반영 답변", "☐"],
    ["6", "HIGH", '구매 의도 감지 ("교육 구매하고 싶다")', "AI가 구매 의도 인식, RFP 작성 제안", "☐"],
    ["7", "HIGH", 'RFP 동의 ("네, 작성해 주세요")', "RFP 유형 선택 화면 표시", "☐"],
    ["8", "HIGH", "RFP 유형 선택 (9개)", "카드 표시, 추천 뱃지, 선택 시 패널 표시", "☐"],
    ["9", "HIGH", '대화로 필드 자동 입력', "RFP 패널에 자동 값 입력", "☐"],
    ["10", "MID", "진행률 바 표시", "필드 입력 시 실시간 갱신", "☐"],
    ["11", "HIGH", "RFP 완료 → Complete 전환", "필수 필드 100% 시 Complete 전환", "☐"],
    ["12", "HIGH", "PDF 다운로드", "나라장터 스타일 PDF 생성, 한글 정상", "☐"],
    ["13", "MID", "공급업체 추천 표시", "카테고리 매칭 업체 목록 + 점수 표시", "☐"],
    ["14", "HIGH", "AI 가드레일: 업체 추천 거부", "특정 업체 추천 요청 시 공정 안내", "☐"],
    ["15", "HIGH", "AI 가드레일: 법적 판단 거부", "법적 판단 요구 시 전문가 상담 권고", "☐"],
    ["16", "HIGH", "환각 방지 (없는 정보)", "지식베이스에 없는 내용 생성 금지", "☐"],
    ["17", "MID", "출처 정확성", "표시된 출처가 실제 답변 내용과 관련", "☐"],
]
add_table(slide, qa1_rows, Inches(0.2), Inches(1.25), Inches(12.9),
          [Inches(0.4), Inches(0.6), Inches(3.0), Inches(4.5), Inches(0.5)], row_height=Inches(0.33))

# ═══════════════════════════════════════════════════
# SLIDE 12 — QA 테스트 체크리스트 (2)
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "10  QA 테스트 체크리스트 (2/2)", "관리자 포털 + 비기능 테스트")
add_footer(slide)

qa2_rows = [
    ["#", "우선", "테스트 항목", "기대 결과", "P/F"],
    ["18", "HIGH", "관리자 정상 로그인", "올바른 계정 입력 시 대시보드 진입", "☐"],
    ["19", "HIGH", "잘못된 비밀번호 입력", '오류 메시지 "비밀번호가 일치하지 않습니다"', "☐"],
    ["20", "MID", "대시보드 통계 표시", "4개 카드 숫자 정상 표시, 1초 이내 로드", "☐"],
    ["21", "HIGH", "PDF 업로드 (지식베이스)", "업로드 → 자동 인제스션 → 성공 메시지", "☐"],
    ["22", "MID", "지식베이스 페이지네이션", "10건 단위 페이징 정상 동작", "☐"],
    ["23", "MID", "헌법 규칙 추가", "유형 선택 → 내용 입력 → 저장 → 목록 반영", "☐"],
    ["24", "MID", "헌법 규칙 삭제", "삭제 버튼 → 확인 모달 → 삭제 완료", "☐"],
    ["25", "MID", "대화 이력 조회", "세션별 대화 목록 + 상세 보기", "☐"],
    ["26", "MID", "공급업체 등록", "필수 정보 입력 → 저장 → 목록 반영", "☐"],
    ["27", "MID", "공급업체 CSV 일괄 등록", "CSV 업로드 → N건 등록 성공", "☐"],
    ["28", "HIGH", "사용자 추가", "이름/이메일/비밀번호/역할 → 생성 완료", "☐"],
    ["29", "HIGH", "중복 이메일 방지", "이미 등록된 이메일 → 오류 메시지", "☐"],
    ["30", "HIGH", "HTTPS 통신 확인", "모든 통신 HTTPS (자물쇠 아이콘)", "☐"],
    ["31", "HIGH", "채팅 P95 응답 시간", "첫 토큰까지 3초 이내", "☐"],
    ["32", "MID", "Chrome / Edge 호환", "전 기능 정상 동작", "☐"],
    ["33", "MID", "모바일 반응형", "모바일 브라우저 레이아웃 정상", "☐"],
    ["34", "LOW", "PDF 인쇄 확인", "A4 사이즈 정상 인쇄", "☐"],
]
add_table(slide, qa2_rows, Inches(0.2), Inches(1.25), Inches(12.9),
          [Inches(0.4), Inches(0.6), Inches(3.0), Inches(4.5), Inches(0.5)], row_height=Inches(0.33))

# 결과 요약 + 서명란
add_text_box(slide, Inches(0.5), Inches(7.0), Inches(5), Inches(0.3),
             "판정: PASS ☑ / FAIL ☒ / N/A —    |    우선순위: HIGH 필수통과 / MID 출시전수정 / LOW 차기업데이트",
             font_size=8, color=GRAY)

# ═══════════════════════════════════════════════════
# SLIDE 13 — 테스트 결과 요약 + 서명
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "테스트 결과 요약 및 승인", "QA Test Summary & Sign-off")
add_footer(slide)

summary_rows = [
    ["테스트 영역", "항목 수", "PASS", "FAIL", "N/A", "통과율"],
    ["챗봇 기본 기능 (접속/채팅/스트리밍)", "5", "", "", "", ""],
    ["페이즈 전환 (Chat→Asked→Filling→Complete)", "6", "", "", "", ""],
    ["RFP 작성 + PDF 다운로드", "3", "", "", "", ""],
    ["공급업체 추천", "2", "", "", "", ""],
    ["AI 응답 품질 + 가드레일", "4", "", "", "", ""],
    ["관리자 로그인 / 인증", "2", "", "", "", ""],
    ["관리자 기능 (지식베이스/헌법/이력/공급업체/사용자)", "10", "", "", "", ""],
    ["비기능 (성능/보안/호환성)", "5", "", "", "", ""],
    ["합계", "34", "", "", "", ""],
]
add_table(slide, summary_rows, Inches(0.5), Inches(1.4), Inches(12.3),
          [Inches(4.5), Inches(1.0), Inches(1.0), Inches(1.0), Inches(1.0), Inches(1.5)], row_height=Inches(0.40))

# 서명란
add_text_box(slide, Inches(0.5), Inches(5.8), Inches(5), Inches(0.3),
             "테스트 승인", font_size=14, color=DARK, bold=True)
sign_rows = [
    ["구분", "성명", "서명", "일자"],
    ["테스트 수행자", "", "", "2026.    .    ."],
    ["검수 담당자", "", "", "2026.    .    ."],
    ["최종 승인자", "", "", "2026.    .    ."],
]
add_table(slide, sign_rows, Inches(0.5), Inches(6.2), Inches(8),
          [Inches(1.8), Inches(2.2), Inches(2.0), Inches(2.0)], row_height=Inches(0.42))

# ═══════════════════════════════════════════════════
# SLIDE 14 — 테스트 가능 카테고리 (지식베이스 현황)
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "테스트 가능 카테고리 (지식베이스 현황)", "현재 청킹 완료된 18개 카테고리 — 전체 분류체계의 약 10%")
add_footer(slide)

add_text_box(slide, Inches(0.5), Inches(1.25), Inches(12), Inches(0.3),
             "아래 18개 카테고리에 대해 RAG 기반 상담 및 RFP 생성 테스트가 가능합니다. (전체 분류체계의 약 10% 수준)",
             font_size=11, color=GRAY)

cat_rows = [
    ["#", "카테고리", "청크 수", "FAQ 수", "문서 수", "테스트 범위"],
    ["1", "OA기기 렌탈 및 유지보수 서비스", "51", "102", "5", "상담 + RFP + 공급업체"],
    ["2", "PC·노트북·모니터·워크스테이션 구매·리스", "22", "43", "3", "상담 + RFP + 공급업체"],
    ["3", "건강검진 서비스", "32", "63", "3", "상담 + RFP + 공급업체"],
    ["4", "계측기 구매·리스", "52", "104", "5", "상담 + RFP + 공급업체"],
    ["5", "공기청정기 렌탈 서비스", "42", "84", "5", "상담 + RFP + 공급업체"],
    ["6", "공장용 소모품", "39", "78", "3", "상담 + RFP + 공급업체"],
    ["7", "법무·노무 자문 서비스", "34", "68", "3", "상담 + RFP + 공급업체"],
    ["8", "복합기·프린터 유지보수 서비스", "54", "106", "5", "상담 + RFP + 공급업체"],
    ["9", "사무가구 구매 및 유지 서비스", "40", "80", "5", "상담 + RFP + 공급업체"],
    ["10", "사무실 인테리어 서비스", "53", "106", "5", "상담 + RFP + 공급업체"],
    ["11", "생산장비 구매·수리 서비스", "50", "100", "5", "상담 + RFP + 공급업체"],
    ["12", "운전용역 서비스", "20", "40", "3", "상담 + RFP + 공급업체"],
    ["13", "의약품 구매 서비스", "33", "26", "4", "상담 + RFP + 공급업체"],
    ["14", "임원차량 리스 서비스", "56", "—", "5", "상담 + RFP"],
    ["15", "청소·소독 서비스", "53", "—", "5", "상담 + RFP"],
    ["16", "칸막이(파티션) 레이아웃 공사", "66", "—", "5", "상담 + RFP"],
    ["17", "택배 서비스", "37", "—", "3", "상담 + RFP"],
    ["18", "회계·세무 용역 서비스", "48", "—", "5", "상담 + RFP"],
]
cat_table = add_table(slide, cat_rows, Inches(0.3), Inches(1.6), Inches(12.7),
          [Inches(0.4), Inches(3.8), Inches(0.9), Inches(0.9), Inches(0.9), Inches(2.0)], row_height=Inches(0.27))

# 합계 행
add_shape(slide, Inches(0.3), Inches(6.95), Inches(12.7), Inches(0.4),
          fill_color=TEAL, border_color=TEAL,
          text="합계:  782 청크  |  1,000 FAQ  |  77 문서  |  18 카테고리",
          font_size=12, font_color=WHITE, bold=True, shape_type=MSO_SHAPE.RECTANGLE)

# ═══════════════════════════════════════════════════
# 저장
# ═══════════════════════════════════════════════════
output_path = os.path.join(os.path.dirname(__file__), "IP_Assist_시스템문서_v1.1.pptx")
prs.save(output_path)
print("PPTX created: " + output_path)
print("Total slides: " + str(len(prs.slides)))
