"""IP_Assist_시스템문서_v1.1.pptx 업데이트 스크립트
v1.0 → v1.1 변경사항 반영
"""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import copy

prs = Presentation('docs/IP_Assist_시스템문서_v1.1.pptx')

# ============================================================
# Helper: find shape containing specific text
# ============================================================
def find_shape_with_text(slide, text):
    for shape in slide.shapes:
        if shape.has_text_frame:
            full = shape.text_frame.text
            if text in full:
                return shape
    return None

def find_table(slide):
    for shape in slide.shapes:
        if shape.has_table:
            return shape.table
    return None

def find_all_tables(slide):
    tables = []
    for shape in slide.shapes:
        if shape.has_table:
            tables.append(shape.table)
    return tables

def set_cell_text(cell, text, size=Pt(9), bold=False, color=None):
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    if color:
        run.font.color.rgb = color

# ============================================================
# Slide 1: 버전 업데이트 v1.0 → v1.1, 날짜 업데이트
# ============================================================
slide1 = prs.slides[0]
for shape in slide1.shapes:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if "v1.0" in run.text:
                    run.text = run.text.replace("v1.0", "v1.1")
                if "2026. 03. 08" in run.text:
                    run.text = run.text.replace("2026. 03. 08", "2026. 03. 10")

print("✓ Slide 1: 버전 v1.1, 날짜 03.10 업데이트")

# ============================================================
# Slide 2: 목차 — 항목 추가 없음 (기존 10개 유지)
# ============================================================
print("✓ Slide 2: 목차 유지")

# ============================================================
# Slide 5 (멀티에이전트): Prefetcher 제거 → 실제 7개만 (8→7행)
# 에이전트 테이블에서 Row 8 (Prefetcher) 텍스트를 비움
# ============================================================
slide5 = prs.slides[4]
tables5 = find_all_tables(slide5)
if tables5:
    agent_table = tables5[0]  # 첫 번째 테이블이 에이전트 테이블
    # Row 8이 Prefetcher — 제거 대신 텍스트를 "미사용"으로 변경
    # (python-pptx는 행 삭제가 어려움)
    last_row = agent_table.rows[8]
    for cell in last_row.cells:
        cell.text = ""
    # Row 7 (RFP Extractor)의 역할 업데이트
    rfp_row = agent_table.rows[7]
    set_cell_text(rfp_row.cells[3], "대화에서 RFP 필드 자동 추출 + 의도 감지(field_input/question/rfp_question)", Pt(8))

print("✓ Slide 5: Prefetcher 행 비움, RFP 에이전트 역할 업데이트")

# ============================================================
# Slide 5: Classification 에이전트 역할에 렌탈 키워드 오버라이드 추가
# ============================================================
if tables5:
    agent_table = tables5[0]
    cls_row = agent_table.rows[2]  # Classification
    set_cell_text(cls_row.cells[3],
        "의도 분류 (구매/질문/RFP동의/필드입력), RFP 유형 추천, 렌탈/리스 키워드 오버라이드", Pt(8))

print("✓ Slide 5: Classification 에이전트에 렌탈 키워드 오버라이드 반영")

# ============================================================
# Slide 7 (DB 설계): conversations 테이블 설명 업데이트
# ============================================================
slide7 = prs.slides[6]
tables7 = find_all_tables(slide7)
if tables7:
    db_table = tables7[0]
    # Row 3 = conversations
    conv_row = db_table.rows[3]
    set_cell_text(conv_row.cells[1],
        "session_id, messages(jsonb), category, rag_score, status", Pt(8))
    set_cell_text(conv_row.cells[2], "B-tree (session_id)", Pt(8))
    set_cell_text(conv_row.cells[3], "대화 이력 자동 저장 (세션별 upsert)", Pt(8))

    # Row 7 = rfp_requests 업데이트
    rfp_row = db_table.rows[7]
    set_cell_text(rfp_row.cells[1],
        "session_id, rfp_type, title, org_name, requester, fields(jsonb), status", Pt(8))
    set_cell_text(rfp_row.cells[3], "RFP 신청 관리 (섹션별 필드 그룹핑)", Pt(8))

print("✓ Slide 7: conversations, rfp_requests 테이블 설명 업데이트")

# ============================================================
# Slide 8 (챗봇 가이드): 후속 질문 설명 업데이트 — 클릭 가능 버튼
# ============================================================
slide8 = prs.slides[7]
shape_suggestion = find_shape_with_text(slide8, "후속 질문")
if shape_suggestion:
    for para in shape_suggestion.text_frame.paragraphs:
        for run in para.runs:
            if "후속 질문" in run.text and "제안" in run.text:
                run.text = "답변 후 3~5개 후속 질문 클릭 버튼 표시. 클릭 시 해당 질문 즉시 전송."

print("✓ Slide 8: 후속 질문 → 클릭 버튼으로 설명 업데이트")

# ============================================================
# Slide 9 (RFP 작성 & 공급업체): PDF 다운로드 설명 업데이트
# ============================================================
slide9 = prs.slides[8]
pdf_shape = find_shape_with_text(slide9, "PDF 다운로드")
if pdf_shape:
    for para in pdf_shape.text_frame.paragraphs:
        for run in para.runs:
            if "나라장터 스타일" in run.text:
                run.text = "• 백엔드 통합 렌더링 (이메일 링크 · PDF 다운로드 동일 로직)"
            elif "구성:" in run.text:
                run.text = "• 구성: 사업명, 발주기관, 담당자, 기간, 목표, 내용, 평가기준, 서명란"
            elif "필수 필드 100%" in run.text:
                run.text = "• 필수 필드 100% 입력 시 다운로드 · 이메일 전송 가능"

print("✓ Slide 9: PDF 다운로드 — 백엔드 통합 렌더링 설명")

# ============================================================
# Slide 10 (관리자 포털): RFP 신청 관리 업데이트
# ============================================================
slide10 = prs.slides[9]
tables10 = find_all_tables(slide10)
if tables10:
    admin_table = tables10[0]
    # Row 4 = 대화 이력
    conv_admin_row = admin_table.rows[4]
    set_cell_text(conv_admin_row.cells[2],
        "세션별 대화 자동 저장. 카테고리, RAG 점수, 상태 필터링. 세션 클릭 시 전체 대화 표시.", Pt(8))

    # Row 5 = RFP 신청
    rfp_admin_row = admin_table.rows[5]
    set_cell_text(rfp_admin_row.cells[2],
        "RFP 목록 조회 (제목·기관·신청자 표시), 섹션별 필드 그룹핑 상세보기, 상태 변경(제출/승인/반려)", Pt(8))

print("✓ Slide 10: 관리자 대화이력/RFP 신청 설명 업데이트")

# ============================================================
# Slide 11 (QA 체크리스트 1/2): 항목 4 업데이트
# ============================================================
slide11 = prs.slides[10]
tables11 = find_all_tables(slide11)
if tables11:
    qa_table = tables11[0]
    # Row 4 = #4 후속 질문 제안 버튼 → 클릭 버튼
    row4 = qa_table.rows[4]
    set_cell_text(row4.cells[2], "후속 질문 클릭 버튼", Pt(8))
    set_cell_text(row4.cells[3], "답변 후 3~5개 버튼 표시, 클릭 시 즉시 전송", Pt(8))

    # Row 8 = #8 RFP 유형 선택 → 클릭 감도 개선
    row8 = qa_table.rows[8]
    set_cell_text(row8.cells[3], "카드 표시, 추천 뱃지, 1회 클릭으로 즉시 선택", Pt(8))

    # Row 12 = #12 PDF 다운로드 — 통합 렌더링
    row12 = qa_table.rows[12]
    set_cell_text(row12.cells[3], "백엔드 통합 렌더링 PDF, 이메일 링크와 동일 출력", Pt(8))

print("✓ Slide 11: QA 체크리스트 항목 4, 8, 12 업데이트")

# ============================================================
# 새 슬라이드: v1.1 변경 이력 (마지막 슬라이드 앞에 삽입)
# ============================================================
# 기존 마지막 슬라이드 뒤에 빈 슬라이드 추가
blank_layout = prs.slide_layouts[6]  # Blank layout
changelog_slide = prs.slides.add_slide(blank_layout)

# 제목
from pptx.util import Inches, Pt, Emu
title_box = changelog_slide.shapes.add_textbox(
    Emu(457200), Emu(274638), Emu(8229600), Emu(548640)
)
tf = title_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
run.text = "v1.1 변경 이력  (2026.03.10)"
run.font.size = Pt(24)
run.font.bold = True
run.font.color.rgb = RGBColor(0x1A, 0x23, 0x7E)

# 부제목
sub_box = changelog_slide.shapes.add_textbox(
    Emu(457200), Emu(823000), Emu(8229600), Emu(365760)
)
tf2 = sub_box.text_frame
p2 = tf2.paragraphs[0]
run2 = p2.add_run()
run2.text = "Change Log — v1.0 → v1.1"
run2.font.size = Pt(14)
run2.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

# 변경 내역 테이블
rows, cols = 9, 3
tbl_shape = changelog_slide.shapes.add_table(
    rows, cols, Emu(457200), Emu(1371600), Emu(8229600), Emu(4572000)
)
tbl = tbl_shape.table

# 열 너비
tbl.columns[0].width = Emu(457200)   # #
tbl.columns[1].width = Emu(2286000)  # 영역
tbl.columns[2].width = Emu(5486400)  # 변경 내용

# 헤더
header_data = ["#", "영역", "변경 내용"]
header_color = RGBColor(0x1A, 0x23, 0x7E)
for i, text in enumerate(header_data):
    cell = tbl.cell(0, i)
    set_cell_text(cell, text, Pt(10), bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    cell.fill.solid()
    cell.fill.fore_color.rgb = header_color

# 변경 내역 데이터
changes = [
    ("1", "대화 이력 저장", "chat/stream, /chat 양쪽 엔드포인트에서 세션 기반 conversations 테이블 자동 upsert. 관리자 대화이력 탭 정상 표시."),
    ("2", "RFP PDF 통합 렌더링", "이메일 링크 RFP 보기(rfp_view.py)와 PDF 다운로드를 동일 백엔드 렌더링 로직으로 통합. 섹션→필드 매핑을 순차 카운트 기반으로 변경."),
    ("3", "RFP 분류 정확도 개선", "렌탈/리스 키워드 오버라이드 추가. '공기청정기 렌탈' → rental 유형 정확 매핑. 비품/소모품·건물관리 등 9개 카테고리 매핑 전면 재점검."),
    ("4", "RFP 질문 맥락 인식", "filling 중 rfp_question 의도 감지 시 사업명+RFP유형으로 검색 쿼리 보강. 해당 카테고리 문서 우선 검색 후 답변 생성."),
    ("5", "후속 질문 클릭 버튼", "제안 텍스트를 클릭 가능 버튼으로 변경. 1-click으로 해당 질문 즉시 전송."),
    ("6", "RFP 유형 선택 UX 개선", "translateY 애니메이션 제거, transition 범위 축소, stopPropagation 적용. 1회 클릭으로 즉시 선택."),
    ("7", "관리자 RFP 상세 보기", "RFP 상세 모달: 제목·기관·신청자 표시, 섹션별 필드 그룹핑 + 색상 구분 dot, 9개 유형별 섹션 구조 정의."),
    ("8", "관리자 RFP 목록 수정", "API select에 title, org_name, requester, rfp_type, fields 컬럼 추가. '제목 없음' 문제 해결."),
]

for idx, (num, area, detail) in enumerate(changes):
    row_idx = idx + 1
    set_cell_text(tbl.cell(row_idx, 0), num, Pt(9))
    set_cell_text(tbl.cell(row_idx, 1), area, Pt(9), bold=True)
    set_cell_text(tbl.cell(row_idx, 2), detail, Pt(8))
    # 줄무늬
    if row_idx % 2 == 0:
        for c in range(cols):
            tbl.cell(row_idx, c).fill.solid()
            tbl.cell(row_idx, c).fill.fore_color.rgb = RGBColor(0xF0, 0xF0, 0xF8)

# 푸터
footer_box = changelog_slide.shapes.add_textbox(
    Emu(457200), Emu(6400800), Emu(8229600), Emu(274638)
)
tf3 = footer_box.text_frame
p3 = tf3.paragraphs[0]
run3 = p3.add_run()
run3.text = "업무마켓9 — IP Assist | (주)캐스팅엔 | 2026.03"
run3.font.size = Pt(8)
run3.font.color.rgb = RGBColor(0x99, 0x99, 0x99)

print("✓ 새 슬라이드: v1.1 변경 이력 추가")

# ============================================================
# 저장
# ============================================================
output_path = 'docs/IP_Assist_시스템문서_v1.1.pptx'
prs.save(output_path)
print(f"\n✅ 저장 완료: {output_path}")
print(f"총 슬라이드: {len(prs.slides)}장")
